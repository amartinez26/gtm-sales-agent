"""
smart_ingest.py
---------------
Incremental, multi-format document ingestion into ChromaDB.

Features:
  - Manifest tracking: skips unchanged files on re-runs
  - Multi-format: .txt .md .pdf .docx .xlsx .xls .csv .pptx
  - Injection scanner: quarantines files with prompt-injection patterns
  - Per-file error isolation: one bad file never stops the batch
  - Full audit log: ingest_errors.log

Usage:
  python smart_ingest.py                        # scans WATCH_PATH from .env
  python smart_ingest.py "\\\\server\\share"    # scans a specific path
"""

import os
import json
import re
import uuid
import shutil
import logging
from pathlib import Path
from typing import Optional

from dotenv import load_dotenv
load_dotenv()

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma

# ── Config ────────────────────────────────────────────────────────────────────
WATCH_PATH      = os.getenv("WATCH_PATH", "./data/stihl_intelligence")
CHROMA_DIR      = "./chroma_db"
MANIFEST_PATH   = os.path.join(CHROMA_DIR, "manifest.json")
ERROR_LOG       = "./ingest_errors.log"
QUARANTINE_DIR  = "./quarantine"
CHUNK_SIZE      = 1000
CHUNK_OVERLAP   = 100
EMBEDDING_MODEL = "all-MiniLM-L6-v2"

SUPPORTED_EXTENSIONS = {'.txt', '.pdf', '.docx', '.xlsx', '.xls', '.csv', '.md', '.pptx'}

# Patterns that indicate prompt injection attempts
INJECTION_PATTERNS = [
    r"ignore\s+(all\s+)?previous\s+instructions",
    r"you\s+are\s+now\s+a",
    r"forget\s+your\s+(previous\s+)?(instructions|training|rules|prompt)",
    r"disregard\s+(all\s+)?(previous\s+)?instructions",
    r"act\s+as\s+(a\s+)?(different|new|another|jailbreak)",
    r"reveal\s+(the\s+)?(system\s+prompt|user\s+quer|internal|api\s+key)",
    r"repeat\s+(all\s+)?(previous\s+|user\s+)?quer",
    r"new\s+instructions?\s*:",
    r"system\s*:\s*(you\s+are|ignore|forget)",
    r"<\s*system\s*>",
    r"\[system\]",
    r"###\s*instruction",
    r"do\s+anything\s+now",
    r"jailbreak",
    r"pretend\s+(you\s+are|to\s+be)",
]

# ── Setup ─────────────────────────────────────────────────────────────────────
os.makedirs(CHROMA_DIR, exist_ok=True)
os.makedirs(QUARANTINE_DIR, exist_ok=True)

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[
        logging.FileHandler(ERROR_LOG),
        logging.StreamHandler(),
    ],
)
log = logging.getLogger(__name__)

# ── Singletons (initialized once per process) ─────────────────────────────────
embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)
vector_db  = Chroma(persist_directory=CHROMA_DIR, embedding_function=embeddings)
splitter   = RecursiveCharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)


# ── Manifest helpers ──────────────────────────────────────────────────────────
def load_manifest() -> dict:
    if os.path.exists(MANIFEST_PATH):
        with open(MANIFEST_PATH, "r") as f:
            return json.load(f)
    return {}


def save_manifest(manifest: dict):
    with open(MANIFEST_PATH, "w") as f:
        json.dump(manifest, f, indent=2)


# ── Injection scanner ─────────────────────────────────────────────────────────
def scan_for_injection(text: str, filepath: str) -> bool:
    """
    Returns True (and quarantines the file) if injection patterns are detected.
    """
    for pattern in INJECTION_PATTERNS:
        if re.search(pattern, text, re.IGNORECASE):
            log.warning(f"[QUARANTINE] Injection pattern '{pattern}' detected in: {filepath}")
            dest = os.path.join(QUARANTINE_DIR, Path(filepath).name + ".quarantine")
            try:
                shutil.copy2(filepath, dest)
                log.warning(f"[QUARANTINE] File copied to: {dest}")
            except Exception as e:
                log.error(f"[QUARANTINE] Could not copy file: {e}")
            return True
    return False


# ── File loaders ──────────────────────────────────────────────────────────────
def load_file(filepath: str) -> Optional[list]:
    """
    Load a file into a list of LangChain Documents.
    Returns None to signal the file should be skipped (logged internally).
    """
    ext = Path(filepath).suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        log.info(f"[SKIP]  {filepath} — unsupported extension '{ext}'")
        return None

    try:
        # ── Plain text / Markdown ──────────────────────────────────────────
        if ext in ('.txt', '.md'):
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            docs = [Document(page_content=text, metadata={"source": filepath})]

        # ── PDF ───────────────────────────────────────────────────────────
        elif ext == '.pdf':
            from langchain_community.document_loaders import PyPDFLoader
            loader = PyPDFLoader(filepath)
            docs = loader.load()

        # ── Word documents ────────────────────────────────────────────────
        elif ext == '.docx':
            from langchain_community.document_loaders import Docx2txtLoader
            loader = Docx2txtLoader(filepath)
            docs = loader.load()

        # ── Excel ─────────────────────────────────────────────────────────
        elif ext in ('.xlsx', '.xls'):
            import pandas as pd
            xl = pd.ExcelFile(filepath)
            all_text = []
            for sheet in xl.sheet_names:
                df = xl.parse(sheet)
                all_text.append(f"--- Sheet: {sheet} ---\n{df.to_string(index=False)}")
            text = "\n\n".join(all_text)
            docs = [Document(page_content=text, metadata={"source": filepath})]

        # ── CSV ───────────────────────────────────────────────────────────
        elif ext == '.csv':
            from langchain_community.document_loaders import CSVLoader
            loader = CSVLoader(filepath)
            docs = loader.load()

        # ── PowerPoint ────────────────────────────────────────────────────
        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(filepath)
            slides_text = []
            for i, slide in enumerate(prs.slides):
                texts = [
                    shape.text for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]
                if texts:
                    slides_text.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))
            text = "\n\n".join(slides_text)
            docs = [Document(page_content=text, metadata={"source": filepath})]

        else:
            log.info(f"[SKIP]  {filepath} — no loader for extension '{ext}'")
            return None

        # Guard: empty content
        if not docs or not any(d.page_content.strip() for d in docs):
            log.info(f"[SKIP]  {filepath} — empty content after loading")
            return None

        return docs

    except PermissionError:
        log.error(f"[ERROR] {filepath} — PermissionError (file locked or access denied)")
        return None
    except Exception as e:
        log.error(f"[ERROR] {filepath} — {type(e).__name__}: {e}")
        return None


# ── Core ingest / remove ──────────────────────────────────────────────────────
def _ingest_file(filepath: str) -> Optional[list]:
    """
    Internal: load, scan, chunk, and add to ChromaDB.
    Returns list of doc_ids on success, None on skip/error.
    """
    docs = load_file(filepath)
    if docs is None:
        return None

    full_text = " ".join(d.page_content for d in docs)
    if scan_for_injection(full_text, filepath):
        return None

    chunks = splitter.split_documents(docs)
    if not chunks:
        log.info(f"[SKIP]  {filepath} — 0 chunks after splitting")
        return None

    doc_ids = [str(uuid.uuid4()) for _ in chunks]
    vector_db.add_documents(documents=chunks, ids=doc_ids)
    log.info(f"[OK]    {filepath} — {len(chunks)} chunks ingested")
    return doc_ids


def _remove_file(filepath: str, manifest: dict):
    """Internal: delete vectors from ChromaDB and remove from manifest."""
    entry = manifest.get(filepath)
    if entry:
        doc_ids = entry.get("doc_ids", [])
        if doc_ids:
            try:
                vector_db.delete(ids=doc_ids)
                log.info(f"[DELETE] {filepath} — {len(doc_ids)} vectors removed")
            except Exception as e:
                log.error(f"[ERROR]  Could not delete vectors for {filepath}: {e}")
        del manifest[filepath]


# ── Public API (used by watchdog_service.py) ──────────────────────────────────
def process_file(filepath: str):
    """
    Called by the watchdog when a file is created or modified.
    Re-ingests the file, replacing any previous vectors for it.
    """
    filepath = str(Path(filepath).resolve())
    manifest = load_manifest()

    if filepath in manifest:
        _remove_file(filepath, manifest)

    try:
        mtime = os.path.getmtime(filepath)
    except OSError:
        return

    doc_ids = _ingest_file(filepath)
    if doc_ids:
        manifest[filepath] = {"last_modified": mtime, "doc_ids": doc_ids}
    save_manifest(manifest)


def delete_file(filepath: str):
    """
    Called by the watchdog when a file is deleted from the watch path.
    """
    filepath = str(Path(filepath).resolve())
    manifest = load_manifest()
    _remove_file(filepath, manifest)
    save_manifest(manifest)


# ── Full / incremental scan ───────────────────────────────────────────────────
def run_full_scan(watch_path: str = WATCH_PATH):
    """
    Walk the entire watch_path and:
      - Skip files unchanged since last run  (manifest timestamp match)
      - Ingest new files
      - Re-ingest modified files
      - Remove vectors for deleted files
    """
    watch_path = str(Path(watch_path).resolve())
    log.info(f"=== Starting scan: {watch_path} ===")
    manifest = load_manifest()

    current_files: set = set()
    new_count = modified_count = skipped_count = error_count = 0

    for root, dirs, files in os.walk(watch_path):
        # Skip hidden directories and common system dirs
        dirs[:] = [
            d for d in dirs
            if not d.startswith('.') and d not in ('$RECYCLE.BIN', 'System Volume Information')
        ]

        for filename in files:
            # Skip hidden files and Office temp files (~$...)
            if filename.startswith('.') or filename.startswith('~$'):
                continue

            filepath = str(Path(os.path.join(root, filename)).resolve())
            ext = Path(filepath).suffix.lower()

            if ext not in SUPPORTED_EXTENSIONS:
                log.info(f"[SKIP]  {filepath} — unsupported extension '{ext}'")
                continue

            current_files.add(filepath)

            try:
                mtime = os.path.getmtime(filepath)
            except OSError:
                continue

            entry = manifest.get(filepath)
            if entry and entry.get("last_modified") == mtime:
                skipped_count += 1
                continue  # unchanged — no work needed

            # New or modified
            if filepath in manifest:
                _remove_file(filepath, manifest)
                modified_count += 1
            else:
                new_count += 1

            doc_ids = _ingest_file(filepath)
            if doc_ids:
                manifest[filepath] = {"last_modified": mtime, "doc_ids": doc_ids}
            else:
                error_count += 1

    # Clean up vectors for files that no longer exist on disk
    deleted_files = set(manifest.keys()) - current_files
    for filepath in deleted_files:
        _remove_file(filepath, manifest)

    save_manifest(manifest)

    log.info(
        f"=== Scan complete — "
        f"new: {new_count} | modified: {modified_count} | "
        f"deleted: {len(deleted_files)} | skipped: {skipped_count} | "
        f"errors/quarantine: {error_count} ==="
    )


# ── Entry point ───────────────────────────────────────────────────────────────
if __name__ == "__main__":
    import sys
    path = sys.argv[1] if len(sys.argv) > 1 else WATCH_PATH
    run_full_scan(path)
